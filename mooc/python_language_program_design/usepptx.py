from pptx import Presentation
from pptx.util import Inches
img_path = 'tree.png'
prs = Presentation()
title_slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(title_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path,left,top)
left = Inches(5)
height = Inches(5.5)
title = slide.shapes.title
pic = slide.shapes.add_picture(img_path,left,top,height=height)

title.text = "hello,world"

prs.save("usepptx.pptx")